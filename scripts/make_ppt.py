# -*- coding: utf-8 -*-
"""生成《交互式四川麻将（血战到底）入门教学》项目汇报 PPT。

用法：python scripts/make_ppt.py
依赖：python-pptx
输出：四川麻将教学_项目汇报.pptx（项目根目录）
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- 主题色 ----------
BRAND = RGBColor(0x0F, 0x51, 0x32)
BRAND2 = RGBColor(0x15, 0x73, 0x47)
ACCENT = RGBColor(0xC8, 0x43, 0x1F)
INK = RGBColor(0x20, 0x28, 0x24)
DIM = RGBColor(0x5D, 0x6B, 0x64)
PANEL = RGBColor(0xF0, 0xF3, 0xEF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xD8, 0xDD, 0xD9)
WAN = RGBColor(0xB3, 0x26, 0x1E)
TONG = RGBColor(0x1F, 0x6F, 0xEB)
TIAO = RGBColor(0x1A, 0x7F, 0x37)

FONT = "Microsoft YaHei"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def _set_font(run, size, color, bold=False, font=FONT):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font
    # 中文字体
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", font)


def box(s, l, t, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE, radius=None):
    sp = s.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def text(s, l, t, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """lines: list of (text, size, color, bold) 或 list of list（多 run 同段）。"""
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(6)
        runs = ln if isinstance(ln, list) else [ln]
        for (txt, size, color, bold) in runs:
            r = p.add_run()
            r.text = txt
            _set_font(r, size, color, bold)
    return tb


def tile(s, l, t, label, color, w=0.62, h=0.86):
    sp = box(s, l, t, w, h, fill=WHITE, line=LINE, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = sp.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    _set_font(r, 15, color, bold=True)
    return sp


def header(s, idx, title, sub=None):
    box(s, 0, 0, SW.inches, 1.15, fill=BRAND)
    box(s, 0, 1.15, SW.inches, 0.06, fill=ACCENT)
    text(s, 0.55, 0.16, 9.5, 0.9,
         [[(f"{idx:02d}  ", 16, RGBColor(0xBF, 0xE3, 0xCF), True), (title, 26, WHITE, True)]],
         anchor=MSO_ANCHOR.MIDDLE)
    if sub:
        text(s, 0.6, 0.74, 11, 0.4, [[(sub, 12, RGBColor(0xD6, 0xE7, 0xDC), False)]])
    text(s, 11.4, 0.16, 1.4, 0.9, [[("川麻教学", 11, RGBColor(0xCF, 0xE3, 0xCF), False)]],
         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def bullets(s, items, l=0.7, t=1.55, w=12.0, size=17, gap=10):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(SH.inches - t - 0.5))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        lvl = 0
        txt = it
        head = None
        if isinstance(it, tuple):
            txt, lvl = it[0], it[1]
            if len(it) > 2:
                head = it[2]
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.level = lvl
        bullet = "▸ " if lvl == 0 else "– "
        if head:
            r = p.add_run(); r.text = bullet + head + "："
            _set_font(r, size, BRAND, True)
            r2 = p.add_run(); r2.text = txt
            _set_font(r2, size, INK, False)
        else:
            r = p.add_run(); r.text = ("    " * lvl) + bullet + txt
            _set_font(r, size - (1 if lvl else 0), INK if lvl == 0 else DIM, lvl == 0)
    return tb


# ---------- 代码页辅助 ----------
CODE_BG = RGBColor(0x15, 0x1B, 0x16)
CODE_FG = RGBColor(0xEA, 0xEF, 0xEA)
MONO = "Consolas"


def code_box(s, l, t, w, h, code, size=11):
    sp = box(s, l, t, w, h, fill=CODE_BG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = sp.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)
    for i, ln in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(0)
        p.line_spacing = 1.05
        r = p.add_run()
        r.text = ln if ln != "" else " "
        r.font.size = Pt(size)
        r.font.name = MONO
        r.font.color.rgb = CODE_FG
    return sp


def code_slide(idx, title, sub, code, notes, code_w=7.15, size=11):
    """左侧代码块 + 右侧逐段说明（notes 为 (要点, 解释) 列表）。"""
    s = slide()
    header(s, idx, title, sub)
    code_box(s, 0.45, 1.45, code_w, 5.55, code, size=size)
    tb = s.shapes.add_textbox(Inches(code_w + 0.7), Inches(1.5),
                              Inches(SW.inches - code_w - 1.15), Inches(5.4))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for head, desc in notes:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(2)
        r = p.add_run(); r.text = "● " + head
        _set_font(r, 13, BRAND, True)
        p2 = tf.add_paragraph()
        p2.space_after = Pt(10)
        r2 = p2.add_run(); r2.text = desc
        _set_font(r2, 11.5, DIM, False)
    return s


# ================= 1. 封面 =================
s = slide()
box(s, 0, 0, SW.inches, SH.inches, fill=BRAND)
box(s, 0, 5.55, SW.inches, 0.05, fill=ACCENT)
text(s, 0.9, 1.7, 11.5, 2.0, [
    [("交互式四川麻将", 46, WHITE, True)],
    [("血战到底 · 入门教学", 30, RGBColor(0xCF, 0xE3, 0xCF), True)],
], anchor=MSO_ANCHOR.TOP)
text(s, 0.95, 3.7, 11, 0.6, [[("纯前端 · 零后端", 16, RGBColor(0xD6, 0xE7, 0xDC), False)]])
# 牌面装饰
demo = [("一万", WAN), ("三条", TIAO), ("五筒", TONG), ("七万", WAN), ("九条", TIAO)]
for i, (lb, c) in enumerate(demo):
    tile(s, 0.95 + i * 0.8, 4.5, lb, c, w=0.66, h=0.9)
text(s, 0.95, 6.0, 11, 0.5, [[("2300012902-舒俊喜", 14, RGBColor(0xCF, 0xE3, 0xCF), False)]])

# ================= 2. 选题背景 =================
s = slide()
header(s, 1, "选题背景：为什么做这个")
bullets(s, [
    ("玩家基数大、但新手最怕“不会算、不敢上桌”。", 0, "真实痛点"),
    ("市面多为纯文字攻略或直接能打的成品游戏；“手把手交互式教规则”几乎是空白。", 0, "差异化"),
    ("牌面渲染、规则演示都不需要后端，可独立完成、好部署。", 0, "纯前端可做"),
    ("规则锁定最普遍认同的“血战到底”，番数采用相加制。", 0, "规则蓝本"),
])

# ================= 3. 项目概览 =================
s = slide()
header(s, 2, "项目概览")
bullets(s, [
    ("原生 HTML + CSS + JavaScript", 0, "形态"),
    ("用 8 关闯关，从认牌、定缺到判胡、算番，零基础学会川麻。", 0, "目标"),
    ("胡牌结构 4 面子 + 1 将 / 七对；必须缺一门；得分 = 底分 × 2^番数。", 0, "规则"),
    ("localStorage 存进度；Service Worker 离线；深色 / 浅色切换。", 0, "工程化"),
], t=1.55, gap=12)
text(s, 0.7, 5.0, 11, 0.5, [[("链接：", 16, INK, False)]])

# ================= 4. 系统架构 =================
s = slide()
header(s, 3, "系统架构：清晰分层")
layers = [
    ("应用层  app.js", "导航 · 进度 · 主题 · Service Worker 注册", BRAND),
    ("关卡层  lessons/ + quiz.js", "8 个关卡 + 通用 createQuiz 闯关框架", BRAND2),
    ("视图层  ui/", "dom.js 构建 · tile.js 牌面 · feedback.js 反馈", RGBColor(0x2E, 0x8B, 0x57)),
    ("引擎层  engine/", "hand · win 胡牌 · ting 听牌 · score 番型算分", ACCENT),
    ("数据层  data/tiles.js", "108 张牌定义 + Unicode 牌面", RGBColor(0x8A, 0x6D, 0x3B)),
]
y = 1.55
for name, desc, col in layers:
    box(s, 1.2, y, 10.9, 1.0, fill=PANEL, line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    box(s, 1.2, y, 0.16, 1.0, fill=col)
    text(s, 1.5, y + 0.08, 5.0, 0.9, [[(name, 17, col, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 6.4, y + 0.08, 5.5, 0.9, [[(desc, 13, DIM, False)]], anchor=MSO_ANCHOR.MIDDLE)
    y += 1.12

# ================= 5. 八关教学内容 =================
s = slide()
header(s, 4, "教学内容：8 关闯关")
lessons = [
    ("1 认牌", "万/条/筒、108 张、无字牌"),
    ("2 定缺", "为何缺一门、花猪惩罚"),
    ("3 基本牌型", "顺子 / 刻子 / 对子，点选组牌"),
    ("4 怎样算胡", "4 面子+1 将，含花猪反例"),
    ("5 碰与杠", "碰 / 明杠 / 暗杠（刮风下雨）"),
    ("6 七对家族", "七对 / 龙七对 / 清七对 / 清龙七对"),
    ("7 番型与算分", "番型拆解 + 底分×2^番数"),
    ("8 综合实战", "13 张听牌：摸哪张能胡"),
]
for i, (t1, t2) in enumerate(lessons):
    col = i % 2
    row = i // 2
    x = 0.9 + col * 6.0
    y = 1.6 + row * 1.32
    box(s, x, y, 5.7, 1.15, fill=PANEL, line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x + 0.25, y + 0.12, 5.2, 0.5, [[(t1, 16, BRAND, True)]])
    text(s, x + 0.25, y + 0.6, 5.2, 0.5, [[(t2, 12, DIM, False)]])

# ================= 5. 项目结构（src/） =================
code_slide(5, "项目结构：src/ 目录一览", "开始读代码前，先看清各文件的分工",
"""src/
├── index.html          入口（空壳）
├── manifest.json       PWA 清单
├── sw.js               Service Worker（离线）
├── icon.svg            应用图标
├── css/style.css       样式（主题/深色/牌面）
├── data/tiles.js       108 张牌 + Unicode
└── js/
    ├── app.js          入口：导航/进度/主题
    ├── engine/         引擎：hand win ting score
    ├── ui/             视图：dom tile feedback
    └── lessons/        关卡：registry quiz gen
        └── lesson-1..8 8 个关卡""",
[("入口三件套", "index.html 空壳 + app.js 启动；manifest/sw 提供 PWA。"),
 ("引擎 engine/", "纯逻辑：胡牌/听牌/算番，可单测、与界面解耦。"),
 ("视图 ui/ + 关卡 lessons/", "ui 负责“画”，lessons 用 quiz 框架“出题”。"),
 ("数据 data/", "牌的编码与 Unicode 牌面，是全项目的地基。")], size=10)


# ================= 代码导读：跟随架构，每层选一个代表 =================

code_slide(6, "代码导读① · 数据层：data/tiles.js", "代表：用一个整数表示一张牌（全项目的地基）",
"""// 花色基数：万=0 / 筒=10 / 条=20；id = base + 点数
export const SUITS = [
  { key:'wan',  name:'万', base:0  },
  { key:'tong', name:'筒', base:10 },
  { key:'tiao', name:'条', base:20 },
];

// 由 id 反查整张牌的信息
export function tileInfo(id) {
  const base  = Math.floor(id / 10) * 10;   // 13 -> 10
  const point = id - base;                  // 13 -> 3
  const suit  = SUITS.find((s) => s.base === base);
  return { id, suit: suit.key, point,
           label: `${CN_NUM[point]}${suit.name}` };  // "三筒"
}""",
[("为何代表数据层", "全项目的“牌”都用这套编码，是一切运算的地基。"),
 ("id = 花色×10 + 点数", "五万=5、三筒=13、九条=29。"),
 ("÷10 取整得花色", "一步从 id 还原出花色基数 0/10/20。"),
 ("tileInfo 反查", "把数字翻译成“三筒”等可读信息供界面用。")])

code_slide(7, "代码导读② · 引擎层：engine/win.js", "代表：胡牌判定（递归+回溯，引擎最难也最亮）",
"""// 剩余牌能否全拆成面子(顺子/刻子)
function canAllMelds(counts) {
  let i = 0;
  while (i < counts.length && counts[i] === 0) i++;
  if (i >= counts.length) return true;     // 全拆完
  if (counts[i] >= 3) {                     // 试刻子
    counts[i] -= 3;
    if (canAllMelds(counts)) { counts[i]+=3; return true; }
    counts[i] += 3;                          // 回溯
  }
  const p = pointOf(i);
  if (p<=7 && counts[i+1]>0 && counts[i+2]>0) {   // 试顺子
    counts[i]--; counts[i+1]--; counts[i+2]--;
    if (canAllMelds(counts)) { counts[i]++;counts[i+1]++;counts[i+2]++; return true; }
    counts[i]++; counts[i+1]++; counts[i+2]++;     // 回溯
  }
  return false;
}
export function isStandardWin(tiles) {
  const counts = toCounts(tiles);
  for (let id=0; id<counts.length; id++)
    if (counts[id] >= 2) {                   // 枚举“将”
      counts[id] -= 2;
      if (canAllMelds(counts)) { counts[id]+=2; return true; }
      counts[id] += 2;
    }
  return false;
}""",
[("为何代表引擎层", "听牌、算番都建立在“能不能胡”之上，是核心。"),
 ("先定将再拆面子", "枚举每种 ≥2 的牌当将，余下交给 canAllMelds。"),
 ("总用最小未用牌 i", "必须消耗它：要么做刻子，要么当顺子头。"),
 ("递归 + 回溯", "试一种拆法失败就把牌放回(counts 复原)再换。")], size=10)

code_slide(8, "代码导读③ · 视图层：ui/dom.js", "代表：el() 一行造出带 class/事件/子节点的元素",
"""// 用法：el('button.choice-btn', { onClick }, '文本')
export function el(tag, props, ...children) {
  let cls = [];
  const tagName = tag.replace(/[.#][^.#]+/g, (m) => {
    if (m[0] === '.') cls.push(m.slice(1)); return '';
  });
  const node = document.createElement(tagName || 'div');
  if (cls.length) node.className = cls.join(' ');
  /* 解析 props：onXxx→事件、dataset、html、普通属性 */
  for (const c of children)
    node.append(c instanceof Node ? c : document.createTextNode(String(c)));
  return node;
}""",
[("为何代表视图层", "整个界面都用它动态生成，没有手写死的 HTML。"),
 ("解析标签", "“.x”进 class、“#y”进 id，剩下当标签名。"),
 ("onXxx → 事件", "以 on 开头的属性自动转成 addEventListener。"),
 ("children → append", "字符串自动转成文本节点挂进去。")])

code_slide(9, "代码导读④ · 关卡层：lessons/lesson-4-win.js", "代表：一关如何用 createQuiz 出题 + 调引擎判分",
"""export default {
  id: 'win', no: 4, title: '怎样算胡',
  build(practiceEl, api) {
    createQuiz(practiceEl, api, {
      target: 4, delay: 700,
      render(ctx) {
        const tiles  = makeCase();                    // 出题
        const answer = canHu(tiles) ? '胡了' : '没胡'; // 引擎给答案
        const options = el('div.quiz-options');
        ['胡了','没胡'].forEach((label) => options.append(
          el('button.choice-btn', { onClick: () =>
            label === answer ? ctx.correct('对！') : ctx.wrong('错') },
            label)));
        ctx.stage.append(
          el('div.hand-board', handEl(tiles, { small:true })), options);
      },
    });
  },
};""",
[("为何代表关卡层", "8 关都是这个套路：出题 → 判定 → ctx.correct/wrong。"),
 ("createQuiz 接管流程", "关卡只写 render(ctx)，计数/过关/停顿交给框架。"),
 ("直接问引擎要答案", "canHu(tiles) 给出标准答案，UI 不关心怎么算。"),
 ("点击对比答案", "选对调 ctx.correct，选错调 ctx.wrong。")], size=10)

code_slide(10, "代码导读⑤ · 应用层：app.js", "代表：把“画布+接口”交给关卡，并管进度/离线",
"""function buildLessonSection() {
  const lesson = LESSONS[currentIndex];
  const practiceEl = el('div.practice');
  const api = {
    feedback(ok, msg){ showFeedback(feedbackEl, ok, msg); },
    clearFeedback(){ clearFeedback(feedbackEl); },
    pass(){ completed.add(lesson.id); saveProgress(); /*刷新进度*/ },
  };
  lesson.build(practiceEl, api);   // 把“画布”和“接口”交给关卡
}

function saveProgress(){           // 进度持久化
  localStorage.setItem(STORAGE_KEY, JSON.stringify([...completed]));
}

// PWA：注册 Service Worker -> 可离线
if ('serviceWorker' in navigator && location.protocol.startsWith('http')) {
  window.addEventListener('load', () =>
    navigator.serviceWorker.register('./sw.js').catch(() => {}));
}""",
[("为何代表应用层", "它是总调度，串起关卡、进度、主题、PWA。"),
 ("api 是关卡契约", "关卡只认 feedback / clearFeedback / pass。"),
 ("localStorage 持久化", "存“已通关 id 集合”，关掉网页也还在。"),
 ("Service Worker", "首访缓存资源，断网也能用(PWA)。")], size=11)


# ================= 技术亮点：亮点介绍 + 结合代码解读 =================

code_slide(11, "技术亮点① · 通用关卡框架 createQuiz", "控制反转：框架管流程，关卡只出题",
"""export function createQuiz(practiceEl, api, { target, delay = 500, render }) {
  let done = 0;                          // 闭包状态：已答对数
  const stage = el('div.quiz-stage');
  practiceEl.append(stage);

  function nextRound() {                  // 出下一轮题
    api.clearFeedback(); clear(stage);
    render(ctx);                          // 让具体关卡渲染一轮
  }
  const ctx = {
    stage,
    correct(msg) {
      done += 1;
      clear(stage);                       // 立即清空 → 防连点
      api.feedback(true, msg);
      if (done >= target) api.pass();     // 够数 → 过关
      else setTimeout(nextRound, delay);  // 否则停顿后下一轮
    },
    wrong(msg) { api.feedback(false, msg); },
  };
  nextRound();
}""",
[("控制反转", "框架管计数/过关/停顿，关卡只实现 render(ctx)。"),
 ("闭包记状态", "done 被内部函数记住，每关独立、不污染全局。"),
 ("防连点", "答对先 clear(stage) 移除按钮，过渡期无可点。"),
 ("价值", "8 关共用一套样板；新增一关只写 render，逻辑收敛。")], size=10)

code_slide(12, "技术亮点② · 番型「相加制」", "用 N 个原子番型，自动覆盖所有组合",
"""// 每识别一个“原子番型”，就加一份 { 名称, 番数 }
if (isPengPengHu(tiles)) detail.push({ name:'对对胡', fan:1 });
if (isQingYiSe(tiles))   detail.push({ name:'清一色', fan:2 });
// …七对(2)、龙(+1)、将对(+3)、断幺九(1) 同理累加

const fans  = detail.reduce((s, d) => s + d.fan, 0);  // 番数相加
const score = base * 2 ** fans;                       // 底分 × 2^番数""",
[("原子化、互不耦合", "每个番型只负责加自己那一份，彼此独立。"),
 ("组合自动成立", "清七对 = 七对(2)+清一色(2)，无需单独写一种。"),
 ("reduce 求和 → 指数计分", "番数相加后 2**fans，每多 1 番收益翻倍。"),
 ("价值", "N 个原子覆盖 2^N 种组合；加新番型只加一行，易维护。")])

code_slide(13, "技术亮点③ · 深色模式：一行换肤", "CSS 变量 + data-theme，改一个属性整页变色",
"""/* 颜色集中成变量；深色块用“同名变量”覆盖 */
:root             { --bg:#f4f6f5; --brand:#0f5132; }  /* 浅色 */
[data-theme=dark] { --bg:#11140f; --brand:#2ea96a; }  /* 深色 */
.lesson { background: var(--bg); }   /* 处处只引用 var() */

// app.js：给 <html> 加/去一个属性即可整页换肤
document.documentElement.setAttribute('data-theme', 'dark');
localStorage.setItem('scmj.theme', 'dark');   // 记住选择""",
[("颜色集中为变量", "全站样式只引用 var(--x)，不写死颜色。"),
 ("同名覆盖", "[data-theme=dark] 与 :root 同专一性，靠书写顺序生效。"),
 ("一处开关 + 记忆", "切换 <html> 的 data-theme，并存 localStorage。"),
 ("价值", "改一个属性，整页换肤，其余样式零改动。")])

code_slide(14, "技术亮点④ · PWA 离线：网络优先", "在线永远最新，断网仍可用",
"""// sw.js：网络优先，失败回退缓存
self.addEventListener('fetch', (e) => {
  e.respondWith(
    fetch(e.request)                        // 先走网络（拿最新）
      .then((res) => { /* 顺便更新缓存 */ return res; })
      .catch(() => caches.match(e.request)) // 失败 → 用缓存（离线）
  );
});

// app.js：特性检测后再注册，老浏览器不报错
if ('serviceWorker' in navigator)
  navigator.serviceWorker.register('./sw.js');""",
[("预缓存应用外壳", "install 时把 html/css/js 先缓存起来。"),
 ("网络优先", "在线总拿最新，改了代码能立刻看到。"),
 ("回退缓存", "断网/弱网时用缓存，实现离线可用。"),
 ("价值", "兼顾“更新及时”与“离线可用”；修了早期缓存优先看不到更新的坑。")])


# ================= AI 使用经验 =================
s = slide()
header(s, 15, "AI 在本项目的使用经验", "全流程协作：选题 → 开发 → 测试 → 文档 → 部署")
bullets(s, [
    ("用 AI 头脑风暴并敲定选题；联网检索、比对、整理权威「血战到底」规则。", 0, "选题与资料"),
    ("AI 协助完成分层架构与全部代码：胡牌/听牌/算番引擎、createQuiz 框架、8 个关卡。", 0, "设计与开发"),
    ("AI 写自测(引擎 19 + 生成器 1204)并做浏览器自动化，主动发现并修复重复点击、SW 缓存等隐蔽 bug。", 0, "测试与纠错"),
    ("AI 自动生成汇报 PPT 与零基础入门导览，并逐条排查 CloudBase 部署报错。", 0, "文档与部署"),
    ("AI 极大加速开发与文档、能主动测试纠错；关键决策仍需人把关——“人定方向、AI 落地、人审校”最有效。", 0, "经验总结"),
], t=1.6, gap=12)


# ================= 10. 总结 =================
s = slide()
box(s, 0, 0, SW.inches, SH.inches, fill=BRAND)
box(s, 0, 2.0, SW.inches, 0.05, fill=ACCENT)
text(s, 0.9, 1.0, 11.5, 1.0, [[("总结与展望", 36, WHITE, True)]])
tb = s.shapes.add_textbox(Inches(0.95), Inches(2.4), Inches(11.4), Inches(4.3))
tf = tb.text_frame; tf.word_wrap = True
points = [
    "一套纯前端、可离线、自带规则引擎的交互式川麻教学。",
    "清晰分层 + createQuiz 通用框架，代码精简、易维护、好汇报。",
    "规则权威、引擎可单测，质量有数据支撑。",
    "后续可拓展：换三张动画、全附加番（天胡/地胡/抢杠）、自由算分计算器、真人对战 AI。",
]
for i, p in enumerate(points):
    par = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    par.space_after = Pt(14)
    r = par.add_run(); r.text = "▸ " + p
    _set_font(r, 18, WHITE, False)
text(s, 0.95, 6.7, 11, 0.5, [[("谢谢观看 · http://localhost:8123/src/", 13, RGBColor(0xCF, 0xE3, 0xCF), False)]])

out = "四川麻将教学_项目汇报.pptx"
try:
    prs.save(out)
except PermissionError:
    import time
    out = f"四川麻将教学_项目汇报_{time.strftime('%H%M%S')}.pptx"
    prs.save(out)
    print("原文件被占用(可能正打开)，已另存为:", out)
print("OK ->", os.path.abspath(out), "slides:", len(prs.slides._sldIdLst))
